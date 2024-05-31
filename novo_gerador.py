# -*- coding: utf-8 -*-


# Importando bibliotecas
import pandas as pd
import math
from pptx import Presentation
import os
import sys
import gspread
from oauth2client.service_account import ServiceAccountCredentials
import tkinter as tk
import webbrowser
from tkinter import messagebox

def celular_amarela():

    # Verifica se estamos executando o código como um script congelado (empacotado pelo PyInstaller)
    if getattr(sys, 'frozen', False):
        # Se estivermos executando como um script congelado, o arquivo .pptx estará no diretório do executável
        pptx_path = os.path.join(sys._MEIPASS, 'celular_amarela.pptx')
        json_path = os.path.join(sys._MEIPASS, 'credentials_serv.json')
    else:
        # Se não estivermos executando como um script congelado, o arquivo .pptx estará no diretório de trabalho atual
        pptx_path = 'celular_amarela.pptx'
        json_path = 'credentials_serv.json'

    # Agora você pode abrir o arquivo .pptx usando o caminho em pptx_path
    with open(pptx_path, 'rb') as f:
        pass  # Faça algo com o arquivo

    # Acessa os dados de uma Planilha Google publicadas em CSV
    path = ('https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUeYcOUktnMQ2hknQF7yO7VFGuCO6APMeWCKXk03Y1Z'
            'ycTIJoAWATed9_PlDoAnnY/pub?gid=1382072451&single=true&output=csv')
    produtos = pd.read_csv(path)

    # Remove os valores nulos
    produtos = produtos.dropna()

    # Alterna os tipos de dados de cada coluna
    produtos['Memória Ram'] = produtos['Memória Ram'].astype(int)
    produtos['Cod Interno'] = produtos['Cod Interno'].astype(int)
    produtos['Preço DE'] = produtos['Preço DE'].str.replace('.', '')
    produtos['Preço DE'] = produtos['Preço DE'].str.replace(',', '.').astype(float)
    produtos['Preço PIX'] = produtos['Preço PIX'].str.replace('.', '')
    produtos['Preço PIX'] = produtos['Preço PIX'].str.replace(',', '.').astype(float)
    produtos['Preço Parcelado'] = produtos['Preço Parcelado'].str.replace('.', '')
    produtos['Preço Parcelado'] = produtos['Preço Parcelado'].str.replace(',', '.').astype(float)
    produtos['Proteção Tela'] = produtos['Proteção Tela'].str.replace('.', '')
    produtos['Proteção Tela'] = produtos['Proteção Tela'].str.replace(',', '.').astype(float)
    produtos['Proteção Total'] = produtos['Proteção Total'].str.replace('.', '')
    produtos['Proteção Total'] = produtos['Proteção Total'].str.replace(',', '.').astype(float)

    # Cria um novo dataframe para armazenar os dados formatados
    produtos_formatados = (
        pd.DataFrame(columns=['Nome do Produto', 'Cod Interno', 'Preço DE', 'Preço PIX', 'Preço Parcelado', 'Câmera',
                              'Processador', 'Memória Ram', 'Tela', 'Proteção Tela', 'Proteção Total']))

    # Convertendo e formatando as colunas conforme especificado
    produtos_formatados['Cod Interno'] = produtos['Cod Interno']
    produtos_formatados['Nome do Produto'] = produtos['Nome do Produto'].str.upper()
    produtos_formatados['Câmera'] = produtos['Câmera'].str.upper()
    produtos_formatados['Processador'] = produtos['Processador'].str.upper()
    produtos_formatados['Tela'] = produtos['Tela'].str.upper()
    produtos_formatados['Memória Ram'] = produtos['Memória Ram'].astype(str) + 'GB DE RAM'
    produtos_formatados['Preço DE'] = ("R$ " + produtos['Preço DE'].apply(
        lambda x: "{:,.2f}".format(x).replace('.', '#').replace(',', '.').replace('#', ',')))
    produtos_formatados['Preço PIX'] = ("R$ " + produtos['Preço PIX'].apply(
        lambda x: "{:,.2f}".format(x).replace('.', '#').replace(',', '.').replace('#', ',')))
    produtos_formatados['Proteção Tela'] = ("+ 12x R$ " + (produtos['Proteção Tela'] / 12).apply(
        lambda x: "{:,.2f}".format(x).replace('.', '#').replace(',', '.').replace('#', ',')))
    produtos_formatados['Proteção Total'] = ("+ 12x R$ " + (produtos['Proteção Total'] / 12).apply(
        lambda x: "{:,.2f}".format(x).replace('.', '#').replace(',', '.').replace('#', ',')))


    def calcular_preco_parcelado(preco):
        if preco > 800:
            parcelas = 10
        else:
            parcelas = math.floor(preco / 80)
        return ("Ou " +
                str(parcelas) +
                "x DE R$ " +
                "{:,.2f}".format(
                    preco / parcelas).replace(
                    '.', '#').replace(',', '.').replace('#', ','))


    produtos_formatados['Preço Parcelado'] = produtos['Preço Parcelado'].apply(calcular_preco_parcelado)
    print('Sucesso!')

    # Carrega a apresentação
    prs = Presentation(pptx_path)  # Aqui está a mudança

    # Adiciona um novo slide com o layout '2_EITQUETA AMARELA'
    slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '2_EITQUETA AMARELA'

    # Para cada linha do dataframe
    for index, row in produtos_formatados.iterrows():
        # Adiciona um novo slide com o layout '2_EITQUETA AMARELA'
        slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '2_EITQUETA AMARELA'
        slide = prs.slides.add_slide(slide_layout)

        # Mapeia cada coluna para um placeholder específico
        placeholders = {10: 'Nome do Produto', 11: 'Cod Interno', 13: 'Preço DE', 14: 'Preço PIX', 15: 'Preço Parcelado',
                        16: 'Câmera', 17: 'Processador', 19: 'Memória Ram', 20: 'Tela', 21: 'Proteção Tela',
                        22: 'Proteção Total'}

        # Para cada placeholder no slide
        for i in placeholders:
            try:
                # Adiciona o texto do dataframe ao placeholder
                slide.placeholders[i].text = str(row[placeholders[i]])
            except KeyError:
                print(f"O slide não tem um placeholder com o índice {i}")

    # Salva a apresentação
    prs.save('nova_celular_amarela.pptx')

    # Carrega os dados salvos das últimas etiquetas geradas e remove os valores nulos
    path = ('https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUeYcOUktnMQ2hknQF7yO7VFGuCO6APMeWCKXk03Y1Zy'
            'cTIJoAWATed9_PlDoAnnY/pub?gid=1893260292&single=true&output=csv')

    df_salvos = pd.read_csv(path)
    df_salvos = df_salvos.dropna()

    # Altera o tipo de dados
    df_salvos['Memória Ram'] = df_salvos['Memória Ram'].astype(int)
    df_salvos['Cod Interno'] = df_salvos['Cod Interno'].astype(int)
    df_salvos['Preço DE'] = df_salvos['Preço DE'].str.replace('.', '')
    df_salvos['Preço DE'] = df_salvos['Preço DE'].str.replace(',', '.').astype(float)
    df_salvos['Preço PIX'] = df_salvos['Preço PIX'].str.replace('.', '')
    df_salvos['Preço PIX'] = df_salvos['Preço PIX'].str.replace(',', '.').astype(float)
    df_salvos['Preço Parcelado'] = df_salvos['Preço Parcelado'].str.replace('.', '')
    df_salvos['Preço Parcelado'] = df_salvos['Preço Parcelado'].str.replace(',', '.').astype(float)
    df_salvos['Proteção Tela'] = df_salvos['Proteção Tela'].str.replace('.', '')
    df_salvos['Proteção Tela'] = df_salvos['Proteção Tela'].str.replace(',', '.').astype(float)
    df_salvos['Proteção Total'] = df_salvos['Proteção Total'].str.replace('.', '')
    df_salvos['Proteção Total'] = df_salvos['Proteção Total'].str.replace(',', '.').astype(float)

    # Primeiro, vamos definir 'Cod Interno' como o índice para ambos os dataframes
    df_salvos.set_index('Cod Interno', inplace=True)
    produtos.set_index('Cod Interno', inplace=True)

    # Agora, vamos remover as linhas em df_salvos que estão em produtos
    df_salvos = df_salvos.loc[~df_salvos.index.isin(produtos.index)]

    # Finalmente, vamos concatenar df_salvos e produtos
    df_final = pd.concat([df_salvos, produtos])

    # Resetando o índice
    df_final.reset_index(inplace=True)

    # Defina o escopo de acesso. Isso especifica quais APIs do Google o código terá acesso.
    scope = ['https://spreadsheets.google.com/feeds', 'https://www.googleapis.com/auth/drive']

    # Use as credenciais do arquivo json para autenticar
    # 'credentials_serv.json' é o arquivo de credenciais da conta de serviço que você baixou do Google Cloud Console.
    creds = ServiceAccountCredentials.from_json_keyfile_name(json_path, scope)

    # Autorize usando as credenciais e inicialize o cliente gspread
    client = gspread.authorize(creds)

    # Abra a planilha do Google usando sua chave única e obtenha a primeira folha
    sheet = client.open_by_key('10emhAinbHEyKEEeQADtFx4yGYe9-JOFrTAzB0_GeLfY').worksheet('Salvos.Celular Amarela')

    # Limpe toda a planilha antes de escrever os dados
    sheet.clear()

    # Atualize a planilha com os dados do DataFrame
    # df_final é o seu DataFrame pandas.
    # Primeiro, convertemos os nomes das colunas do DataFrame em uma lista.
    # Em seguida, convertemos cada linha do DataFrame em uma lista.
    # Finalmente, usamos o método update para escrever essas listas na planilha do Google Sheets.
    sheet.update([df_final.columns.values.tolist()] + df_final.values.tolist())

def celular_branca():


    # Verifica se estamos executando o código como um script congelado (empacotado pelo PyInstaller)
    if getattr(sys, 'frozen', False):
        # Se estivermos executando como um script congelado, o arquivo .pptx estará no diretório do executável
        pptx_path = os.path.join(sys._MEIPASS, 'etiqueta_branca.pptx')
        json_path = os.path.join(sys._MEIPASS, 'credentials_serv.json')
    else:
        # Se não estivermos executando como um script congelado, o arquivo .pptx estará no diretório de trabalho atual
        pptx_path = 'etiqueta_branca.pptx'
        json_path = 'credentials_serv.json'

    # Agora você pode abrir o arquivo .pptx usando o caminho em pptx_path
    with open(pptx_path, 'rb') as f:
        pass  # Faça algo com o arquivo

    # Acessa os dados de uma Planilha Google publicadas em CSV
    path = ('https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUeYcOUktnMQ2hknQF7yO7VFGuCO6APMeWCK'
            'Xk03Y1ZycTIJoAWATed9_PlDoAnnY/pub?gid=602087846&single=true&output=csv')
    produtos = pd.read_csv(path)
    # Remove os valores nulos
    produtos = produtos.dropna()

    # Altera o tipo de dados
    produtos['Memória Ram'] = produtos['Memória Ram'].astype(int)
    produtos['Cod Interno'] = produtos['Cod Interno'].astype(int)
    produtos['Preço Por'] = produtos['Preço Por'].str.replace('.', '')
    produtos['Preço Por'] = produtos['Preço Por'].str.replace(',', '.').astype(float)
    produtos['Proteção Tela'] = produtos['Proteção Tela'].str.replace('.', '')
    produtos['Proteção Tela'] = produtos['Proteção Tela'].str.replace(',', '.').astype(float)
    produtos['Proteção Total'] = produtos['Proteção Total'].str.replace('.', '')
    produtos['Proteção Total'] = produtos['Proteção Total'].str.replace(',', '.').astype(float)

    # Cria um novo dataframe para armazenar os dados formatados
    produtos_formatados = (
        pd.DataFrame(columns=[
            'Cod Interno', 'Nome do Produto', 'Preço Por', 'Preço Parcelado', 'Câmera',
            'Processador', 'Memória Ram', 'Tela', 'Proteção Tela', 'Proteção Total']))

    # Convertendo e formatando as colunas conforme especificado
    produtos_formatados['Cod Interno'] = produtos['Cod Interno']
    produtos_formatados['Nome do Produto'] = produtos['Nome do Produto'].str.upper()
    produtos_formatados['Câmera'] = produtos['Câmera'].str.upper()
    produtos_formatados['Processador'] = produtos['Processador'].str.upper()
    produtos_formatados['Tela'] = produtos['Tela'].str.upper()
    produtos_formatados['Memória Ram'] = produtos['Memória Ram'].astype(str) + 'GB DE RAM'
    produtos_formatados['Preço Por'] = (
            "R$ " + produtos['Preço Por'].apply(lambda x: "{:,.2f}".format(x).replace('.', '#').
                                                replace(',', '.').replace('#', ',')))
    produtos_formatados['Proteção Tela'] = (
            "+ 12x R$ " +
            (produtos['Proteção Tela'] / 12).apply(
                lambda x:"{:,.2f}".format(x).replace(
                    '.', '#').replace(',','.').replace('#', ',')))
    produtos_formatados['Proteção Total'] = (
            "+ 12x R$ " + (produtos['Proteção Total'] / 12).apply(
        lambda x: "{:,.2f}".format(x).replace('.', '#').replace(',', '.').replace('#', ',')))

    def calcular_preco_parcelado(preco):
        if preco > 800:
            parcelas = 10
        else:
            parcelas = math.floor(preco / 80)
        return "Ou " + str(parcelas) + "x DE R$ " + "{:,.2f}".format(preco / parcelas).replace('.', '#').replace(',',
                                                                                                                 '.').replace(
            '#', ',')

    produtos_formatados['Preço Parcelado'] = produtos['Preço Por'].apply(calcular_preco_parcelado)

    print('Sucesso!')

    # Carrega a apresentação
    prs = Presentation(pptx_path)  # Aqui está a mudança

    # Adiciona um novo slide com o layout '2_EITQUETA AMARELA'
    slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '2_EITQUETA AMARELA'

    # Para cada linha do dataframe
    for index, row in produtos_formatados.iterrows():
        # Adiciona um novo slide com o layout '2_EITQUETA AMARELA'
        slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '2_EITQUETA AMARELA'
        slide = prs.slides.add_slide(slide_layout)

        # Mapeia cada coluna para um placeholder específico
        placeholders = {10: 'Nome do Produto', 11: 'Cod Interno', 14: 'Preço Por', 15: 'Preço Parcelado', 16: 'Câmera',
                        17: 'Processador', 19: 'Memória Ram', 20: 'Tela', 21: 'Proteção Tela', 22: 'Proteção Total'}

        # Para cada placeholder no slide
        for i in placeholders:
            try:
                # Adiciona o texto do dataframe ao placeholder
                slide.placeholders[i].text = str(row[placeholders[i]])
            except KeyError:
                print(f"O slide não tem um placeholder com o índice {i}")

    # Salva a apresentação
    prs.save('nova_celular_branca.pptx')

    # vCarrega os dados salvos das últimas etiquetas geradas e remove os valores nulos
    path = 'https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUeYcOUktnMQ2hknQF7yO7VFGuCO6APMeWCKXk03Y1ZycTIJoAWATed9_PlDoAnnY/pub?gid=1774661786&single=true&output=csv'
    df_salvos = pd.read_csv(path)

    # vRemove os valores nulos
    df_salvos = df_salvos.dropna()
    # Altera o tipo de coluna para int
    df_salvos['Memória Ram'] = df_salvos['Memória Ram'].astype(int)
    df_salvos['Cod Interno'] = df_salvos['Cod Interno'].astype(int)

    # vAltera o tipo de coluna para float
    df_salvos['Preço Por'] = df_salvos['Preço Por'].str.replace('.', '')
    df_salvos['Preço Por'] = df_salvos['Preço Por'].str.replace(',', '.').astype(float)

    df_salvos['Proteção Tela'] = df_salvos['Proteção Tela'].str.replace('.', '')
    df_salvos['Proteção Tela'] = df_salvos['Proteção Tela'].str.replace(',', '.').astype(float)

    df_salvos['Proteção Total'] = df_salvos['Proteção Total'].str.replace('.', '')
    df_salvos['Proteção Total'] = df_salvos['Proteção Total'].str.replace(',', '.').astype(float)

    # Primeiro, vamos definir 'Cod Interno' como o índice para ambos os dataframes
    df_salvos.set_index('Cod Interno', inplace=True)
    produtos.set_index('Cod Interno', inplace=True)

    # Agora, vamos remover as linhas em df_salvos que estão em produtos
    df_salvos = df_salvos.loc[~df_salvos.index.isin(produtos.index)]

    # Finalmente, vamos concatenar df_salvos e produtos
    df_final = pd.concat([df_salvos, produtos])

    # Resetando o índice
    df_final.reset_index(inplace=True)

    # Defina o escopo de acesso. Isso especifica quais APIs do Google o código terá acesso.
    scope = ['https://spreadsheets.google.com/feeds', 'https://www.googleapis.com/auth/drive']

    # Use as credenciais do arquivo json para autenticar
    # 'credentials_serv.json' é o arquivo de credenciais da conta de serviço que você baixou do Google Cloud Console.
    creds = ServiceAccountCredentials.from_json_keyfile_name(json_path, scope)

    # Autorize usando as credenciais e inicialize o cliente gspread
    client = gspread.authorize(creds)

    # Abra a planilha do Google usando sua chave única e obtenha a primeira folha
    sheet = client.open_by_key('10emhAinbHEyKEEeQADtFx4yGYe9-JOFrTAzB0_GeLfY').worksheet('Salvos.Celular Branca')

    # Limpe toda a planilha antes de escrever os dados
    sheet.clear()

    # Atualize a planilha com os dados do DataFrame
    # df_final é o seu DataFrame pandas.
    # Primeiro, convertemos os nomes das colunas do DataFrame em uma lista.
    # Em seguida, convertemos cada linha do DataFrame em uma lista.
    # Finalmente, usamos o método update para escrever essas listas na planilha do Google Sheets.
    sheet.update([df_final.columns.values.tolist()] + df_final.values.tolist())

def notebook_amarela():

    # Verifica se estamos executando o código como um script congelado (empacotado pelo PyInstaller)
    if getattr(sys, 'frozen', False):
        # Se estivermos executando como um script congelado, o arquivo .pptx estará no diretório do executável
        pptx_path = os.path.join(sys._MEIPASS, 'notebook_amarela.pptx')
        json_path = os.path.join(sys._MEIPASS, 'credentials_serv.json')
    else:
        # Se não estivermos executando como um script congelado, o arquivo .pptx estará no diretório de trabalho atual
        pptx_path = 'notebook_amarela.pptx'
        json_path = 'credentials_serv.json'

    # Agora você pode abrir o arquivo .pptx usando o caminho em pptx_path
    with open(pptx_path, 'rb') as f:
        pass  # Faça algo com o arquivo

    # Acessa os dados de uma Planilha Google publicadas em CSV
    path = (''
            'https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUeYcOUktnMQ2hknQF7y'
            'O7VFGuCO6APMeWCKXk03Y1ZycTIJoAWATed9_PlDoAnnY/pub?gid=203841021&single=true&output=csv')
    produtos = pd.read_csv(path)

    # Remove os valores nulos
    produtos = produtos.dropna()

    # Altera o tipo de coluna para int
    produtos['Memória Ram'] = produtos['Memória Ram'].astype(int)
    produtos['Cod Interno'] = produtos['Cod Interno'].astype(int)
    produtos['Preço DE'] = produtos['Preço DE'].str.replace('.', '')
    produtos['Preço DE'] = produtos['Preço DE'].str.replace(',', '.').astype(float)
    produtos['Preço PIX'] = produtos['Preço PIX'].str.replace('.', '')
    produtos['Preço PIX'] = produtos['Preço PIX'].str.replace(',', '.').astype(float)
    produtos['Preço Parcelado'] = produtos['Preço Parcelado'].str.replace('.', '')
    produtos['Preço Parcelado'] = produtos['Preço Parcelado'].str.replace(',', '.').astype(float)
    produtos['GE'] = produtos['GE'].str.replace('.', '')
    produtos['GE'] = produtos['GE'].str.replace(',', '.').astype(float)

    # Cria um novo dataframe para armazenar os dados formatados
    produtos_formatados = (
        pd.DataFrame(columns=['Nome do Produto', 'Cod Interno', 'Preço DE', 'Preço PIX',
                              'Preço Parcelado', 'Armazenamento', 'Processador', 'Memória Ram', 'Tela', 'Graficos',
                              'GE']))

    # Convertendo e formatando as colunas conforme especificado
    produtos_formatados['Cod Interno'] = produtos['Cod Interno']
    produtos_formatados['Nome do Produto'] = produtos['Nome do Produto'].str.upper()
    produtos_formatados['Armazenamento'] = produtos['Armazenamento'].str.upper()
    produtos_formatados['Processador'] = produtos['Processador'].str.upper()
    produtos_formatados['Tela'] = produtos['Tela'].str.upper()
    produtos_formatados['Graficos'] = produtos['Graficos'].str.upper()
    produtos_formatados['Memória Ram'] = produtos['Memória Ram'].astype(str) + 'GB DE RAM'
    produtos_formatados['Preço DE'] = (
            "R$ " + produtos['Preço DE'].apply(lambda x: "{:,.2f}".format(x).replace('.', '#').
                                               replace(',', '.').replace('#', ',')))
    produtos_formatados['Preço PIX'] = (
            "R$ " + produtos['Preço PIX'].apply(lambda x: "{:,.2f}".format(x).replace('.', '#').
                                                replace(',', '.').replace('#', ',')))
    produtos_formatados['GE'] = (
            "+ 12x R$ " + (produtos['GE'] / 12).apply(lambda x:
                                                      "{:,.2f}".format(x).replace('.', '#').replace(',', '.').replace
                                                      ('#', ',')))

    def calcular_preco_parcelado(preco):
        if preco > 800:
            parcelas = 10
        else:
            parcelas = math.floor(preco / 80)
        return ("Ou " + str(parcelas) + "x DE R$ " +
                "{:,.2f}".format(preco / parcelas).replace('.', '#').replace(',', '.').replace('#', ','))

    produtos_formatados['Preço Parcelado'] = produtos['Preço Parcelado'].apply(calcular_preco_parcelado)
    print('Sucesso!')

    # Carrega a apresentação
    prs = Presentation(pptx_path)  # Aqui está a mudança

    # Adiciona um novo slide com o layout '
    slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '

    # Para cada linha do dataframe
    for index, row in produtos_formatados.iterrows():
        # Adiciona um novo slide com o layout '
        slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '
        slide = prs.slides.add_slide(slide_layout)

        # Mapeia cada coluna para um placeholder específico
        placeholders = {10: 'Nome do Produto', 11: 'Cod Interno', 13: 'Preço DE', 14: 'Preço PIX',
                        15: 'Preço Parcelado', 16: 'Armazenamento', 17: 'Processador', 19: 'Memória Ram', 20: 'Tela',
                        21: 'Graficos', 22: 'GE'}

        # Para cada placeholder no slide
        for i in placeholders:
            try:
                # Adiciona o texto do dataframe ao placeholder
                slide.placeholders[i].text = str(row[placeholders[i]])
            except KeyError:
                print(f"O slide não tem um placeholder com o índice {i}")

    # Salva a apresentação
    prs.save('nova_notebook_amarela.pptx')

    # Carrega os dados salvos das últimas etiquetas geradas e remove os valores nulos
    path = \
        ('https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUe'
         'YcOUktnMQ2hknQF7yO7VFGuCO6APMeWCKXk03Y1ZycTIJoAWATed9_PlDoAnnY/pub?gid=2105789903&single=true&output=csv')
    df_salvos = pd.read_csv(path)
    # Remove os valores nulos
    df_salvos = df_salvos.dropna()

    # Altera o tipo de coluna para int
    df_salvos['Memória Ram'] = df_salvos['Memória Ram'].astype(int)
    df_salvos['Cod Interno'] = df_salvos['Cod Interno'].astype(int)
    df_salvos['Preço DE'] = df_salvos['Preço DE'].str.replace('.', '')
    df_salvos['Preço DE'] = df_salvos['Preço DE'].str.replace(',', '.').astype(float)
    df_salvos['Preço PIX'] = df_salvos['Preço PIX'].str.replace('.', '')
    df_salvos['Preço PIX'] = df_salvos['Preço PIX'].str.replace(',', '.').astype(float)
    df_salvos['Preço Parcelado'] = df_salvos['Preço Parcelado'].str.replace('.', '')
    df_salvos['Preço Parcelado'] = df_salvos['Preço Parcelado'].str.replace(',', '.').astype(float)
    df_salvos['GE'] = df_salvos['GE'].str.replace('.', '')
    df_salvos['GE'] = df_salvos['GE'].str.replace(',', '.').astype(float)

    # Primeiro, vamos definir 'Cod Interno' como o índice para ambos os dataframes
    df_salvos.set_index('Cod Interno', inplace=True)
    produtos.set_index('Cod Interno', inplace=True)

    # Agora, vamos remover as linhas em df_salvos que estão em produtos
    df_salvos = df_salvos.loc[~df_salvos.index.isin(produtos.index)]

    # Finalmente, vamos concatenar df_salvos e produtos
    df_final = pd.concat([df_salvos, produtos])

    # Resetando o índice
    df_final.reset_index(inplace=True)

    # Defina o escopo de acesso. Isso especifica quais APIs do Google o código terá acesso.
    scope = ['https://spreadsheets.google.com/feeds', 'https://www.googleapis.com/auth/drive']

    # Use as credenciais do arquivo json para autenticar
    # 'credentials_serv.json' é o arquivo de credenciais da conta de serviço que você baixou do Google Cloud Console.
    creds = ServiceAccountCredentials.from_json_keyfile_name(json_path, scope)

    # Autorize usando as credenciais e inicialize o cliente gspread
    client = gspread.authorize(creds)

    # Abra a planilha do Google usando sua chave única e obtenha a primeira folha
    sheet = client.open_by_key('10emhAinbHEyKEEeQADtFx4yGYe9-JOFrTAzB0_GeLfY').worksheet('Salvos.Notebook Amarela')

    # Limpe toda a planilha antes de escrever os dados
    sheet.clear()

    # Atualize a planilha com os dados do DataFrame
    # df_final é o seu DataFrame pandas.
    # Primeiro, convertemos os nomes das colunas do DataFrame em uma lista.
    # Em seguida, convertemos cada linha do DataFrame em uma lista.
    # Finalmente, usamos o método update para escrever essas listas na planilha do Google Sheets.
    sheet.update([df_final.columns.values.tolist()] + df_final.values.tolist())

def notebook_branca():
    # Verifica se estamos executando o código como um script congelado (empacotado pelo PyInstaller)
    if getattr(sys, 'frozen', False):
        # Se estivermos executando como um script congelado, o arquivo .pptx estará no diretório do executável
        pptx_path = os.path.join(sys._MEIPASS, 'notebook_branca.pptx')
        json_path = os.path.join(sys._MEIPASS, 'credentials_serv.json')
    else:
        # Se não estivermos executando como um script congelado, o arquivo .pptx estará no diretório de trabalho atual
        pptx_path = 'notebook_branca.pptx'
        json_path = 'credentials_serv.json'

    # Agora você pode abrir o arquivo .pptx usando o caminho em pptx_path
    with open(pptx_path, 'rb') as f:
        pass  # Faça algo com o arquivo

    # Acessa os dados de uma Planilha Google publicadas em CSV
    path = ('https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUeYcO'
            'UktnMQ2hknQF7yO7VFGuCO6APMeWCKXk03Y1ZycTIJoAWATed9_PlDoAnnY/''pub?gid=169319900&single=true&output=csv')
    produtos = pd.read_csv(path)
    # Remove os valores nulos
    produtos = produtos.dropna()

    # Altera o tipo de dados
    produtos['Memória Ram'] = produtos['Memória Ram'].astype(int)
    produtos['Cod Interno'] = produtos['Cod Interno'].astype(int)
    produtos['Preço Por'] = produtos['Preço Por'].str.replace('.', '')
    produtos['Preço Por'] = produtos['Preço Por'].str.replace(',', '.').astype(float)
    produtos['GE'] = produtos['GE'].str.replace('.', '')
    produtos['GE'] = produtos['GE'].str.replace(',', '.').astype(float)

    # Cria um novo dataframe para armazenar os dados formatados
    produtos_formatados = (
        pd.DataFrame(columns=['Nome do Produto', 'Cod Interno', 'Preço Por', 'Preço Parcelado',
                              'Armazenamento', 'Processador', 'Memória Ram', 'Tela', 'Graficos', 'GE']))

    # Convertendo e formatando as colunas conforme especificado
    produtos_formatados['Cod Interno'] = produtos['Cod Interno']
    produtos_formatados['Nome do Produto'] = produtos['Nome do Produto'].str.upper()
    produtos_formatados['Armazenamento'] = produtos['Armazenamento'].str.upper()
    produtos_formatados['Processador'] = produtos['Processador'].str.upper()
    produtos_formatados['Tela'] = produtos['Tela'].str.upper()
    produtos_formatados['Graficos'] = produtos['Graficos'].str.upper()
    produtos_formatados['Memória Ram'] = produtos['Memória Ram'].astype(str) + 'GB DE RAM'
    produtos_formatados['Preço Por'] = (
            "R$ " + produtos['Preço Por'].apply(lambda x: "{:,.2f}".format(x).replace('.', '#').
                                                replace(',', '.').replace('#', ',')))
    produtos_formatados['GE'] = (
            "+ 12x R$ " + (produtos['GE'] / 12).apply(lambda x:
                                                      "{:,.2f}".format(x).replace('.', '#').replace(',', '.').replace
                                                      ('#', ',')))

    def calcular_preco_parcelado(preco):
        if preco > 800:
            parcelas = 10
        else:
            parcelas = math.floor(preco / 80)
        return ("Ou " + str(parcelas) + "x DE R$ " +
                "{:,.2f}".format(preco / parcelas).replace('.', '#').replace(',', '.').replace('#', ','))

    produtos_formatados['Preço Parcelado'] = produtos['Preço Por'].apply(calcular_preco_parcelado)
    print('Sucesso!')

    # Carrega a apresentação
    prs = Presentation(pptx_path)  # Aqui está a mudança

    # Adiciona um novo slide com o layout '
    slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '

    # Para cada linha do dataframe
    for index, row in produtos_formatados.iterrows():
        # Adiciona um novo slide com o layout '
        slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '
        slide = prs.slides.add_slide(slide_layout)

        # Mapeia cada coluna para um placeholder específico
        placeholders = \
            {10: 'Nome do Produto', 11: 'Cod Interno', 14:
                'Preço Por', 15: 'Preço Parcelado', 16: 'Armazenamento', 17:
                 'Processador', 19: 'Memória Ram', 20: 'Tela', 21: 'Graficos', 22: 'GE'}

        # Para cada placeholder no slide
        for i in placeholders:
            try:
                # Adiciona o texto do dataframe ao placeholder
                slide.placeholders[i].text = str(row[placeholders[i]])
            except KeyError:
                print(f"O slide não tem um placeholder com o índice {i}")

    # Salva a apresentação
    prs.save('nova_notebook_branca.pptx')

    # Carrega os dados salvos das últimas etiquetas geradas e remove os valores nulos
    path = (
        'https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUeYcOUktnMQ2hknQF7yO7V'
        'FGuCO6APMeWCKXk03Y1ZycTIJoAWATed9_PlDoAnnY/pub?gid=1159710822&single=true&output=csv')
    df_salvos = pd.read_csv(path)
    # Remove os valores nulos
    df_salvos = df_salvos.dropna()

    # Altera o tipo de coluna para int
    df_salvos['Memória Ram'] = df_salvos['Memória Ram'].astype(int)
    df_salvos['Cod Interno'] = df_salvos['Cod Interno'].astype(int)
    df_salvos['Preço Por'] = df_salvos['Preço Por'].str.replace('.', '')
    df_salvos['Preço Por'] = df_salvos['Preço Por'].str.replace(',', '.').astype(float)
    df_salvos['GE'] = df_salvos['GE'].str.replace('.', '')
    df_salvos['GE'] = df_salvos['GE'].str.replace(',', '.').astype(float)

    # Primeiro, vamos definir 'Cod Interno' como o índice para ambos os dataframes
    df_salvos.set_index('Cod Interno', inplace=True)
    produtos.set_index('Cod Interno', inplace=True)

    # Agora, vamos remover as linhas em df_salvos que estão em produtos
    df_salvos = df_salvos.loc[~df_salvos.index.isin(produtos.index)]

    # Finalmente, vamos concatenar df_salvos e produtos
    df_final = pd.concat([df_salvos, produtos])

    # Resetando o índice
    df_final.reset_index(inplace=True)

    # Defina o escopo de acesso. Isso especifica quais APIs do Google o código terá acesso.
    scope = ['https://spreadsheets.google.com/feeds', 'https://www.googleapis.com/auth/drive']

    # Use as credenciais do arquivo json para autenticar
    # 'credentials_serv.json' é o arquivo de credenciais da conta de serviço que você baixou do Google Cloud Console.
    creds = ServiceAccountCredentials.from_json_keyfile_name(json_path, scope)

    # Autorize usando as credenciais e inicialize o cliente gspread
    client = gspread.authorize(creds)

    # Abra a planilha do Google usando sua chave única e obtenha a primeira folha
    sheet = client.open_by_key('10emhAinbHEyKEEeQADtFx4yGYe9-JOFrTAzB0_GeLfY').worksheet('Salvos.Notebook Branca')

    # Limpe toda a planilha antes de escrever os dados
    sheet.clear()

    # Atualize a planilha com os dados do DataFrame
    # df_final é o seu DataFrame pandas.
    # Primeiro, convertemos os nomes das colunas do DataFrame em uma lista.
    # Em seguida, convertemos cada linha do DataFrame em uma lista.
    # Finalmente, usamos o método update para escrever essas listas na planilha do Google Sheets.
    sheet.update([df_final.columns.values.tolist()] + df_final.values.tolist())

def relampago():
    # Verifica se estamos executando o código como um script congelado (empacotado pelo PyInstaller)
    if getattr(sys, 'frozen', False):
        # Se estivermos executando como um script congelado, o arquivo .pptx estará no diretório do executável
        pptx_path = os.path.join(sys._MEIPASS, 'promocao_relampago.pptx')
        json_path = os.path.join(sys._MEIPASS, 'credentials_serv.json')
    else:
        # Se não estivermos executando como um script congelado, o arquivo .pptx estará no diretório de trabalho atual
        pptx_path = 'promocao_relampago.pptx'
        json_path = 'credentials_serv.json'

    # Agora você pode abrir o arquivo .pptx usando o caminho em pptx_path
    with open(pptx_path, 'rb') as f:
        pass  # Faça algo com o arquivo

    # Acessa os dados de uma Planilha Google publicadas em CSV
    path = \
        ('https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUeYc'
         'OUktnMQ2hknQF7yO7VFGuCO6APMeWCKXk03Y1ZycTIJoAWATed9_PlDoAnnY/pub?gid=24918558&single=true&output=csv')
    produtos = pd.read_csv(path)
    # Remove os valores nulos
    produtos = produtos.dropna()

    # Altera o tipo de coluna
    produtos['Cod Interno'] = produtos['Cod Interno'].astype(int)
    produtos['Preço DE'] = produtos['Preço DE'].str.replace('.', '')
    produtos['Preço DE'] = produtos['Preço DE'].str.replace(',', '.').astype(float)
    produtos['Preço PIX'] = produtos['Preço PIX'].str.replace('.', '')
    produtos['Preço PIX'] = produtos['Preço PIX'].str.replace(',', '.').astype(float)
    produtos['Preço Parcelado'] = produtos['Preço Parcelado'].str.replace('.', '')
    produtos['Preço Parcelado'] = produtos['Preço Parcelado'].str.replace(',', '.').astype(float)

    # Cria um novo dataframe para armazenar os dados formatados
    produtos_formatados = (
        pd.DataFrame(columns=['Nome do Produto', 'Cod Interno', 'Preço DE', 'Preço PIX', 'Preço Parcelado']))

    # Convertendo e formatando as colunas conforme especificado
    produtos_formatados['Cod Interno'] = produtos['Cod Interno']
    produtos_formatados['Nome do Produto'] = produtos['Nome do Produto'].str.upper()
    produtos_formatados['Preço DE'] = \
        ("R$ " + produtos['Preço DE'].apply(lambda x: "{:,.2f}".format(x).replace('.', '#').
                                            replace(',', '.').replace('#', ',')))
    produtos_formatados['Preço PIX'] = (
            "R$ " + produtos['Preço PIX'].apply(lambda x: "{:,.2f}".format(x).replace('.', '#').
                                                replace(',', '.').replace('#', ',')))

    def calcular_preco_parcelado(preco):
        if preco > 800:
            parcelas = 10
        else:
            parcelas = math.floor(preco / 80)
        return ("Ou " + str(parcelas) + "x DE R$ " +
                "{:,.2f}".format(preco / parcelas).replace('.', '#').replace(',', '.').replace('#', ','))

    produtos_formatados['Preço Parcelado'] = produtos['Preço Parcelado'].apply(calcular_preco_parcelado)
    print('Sucesso!')

    # Carrega a apresentação
    prs = Presentation(pptx_path)  # Aqui está a mudança

    # Adiciona um novo slide com o layout '
    slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '

    # Para cada linha do dataframe
    for index, row in produtos_formatados.iterrows():
        # Adiciona um novo slide com o layout '
        slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '
        slide = prs.slides.add_slide(slide_layout)

        # Mapeia cada coluna para um placeholder específico
        placeholders = {10: 'Nome do Produto', 11: 'Cod Interno', 13: 'Preço DE', 14: 'Preço PIX',
                        15: 'Preço Parcelado'}

        # Para cada placeholder no slide
        for i in placeholders:
            try:
                # Adiciona o texto do dataframe ao placeholder
                slide.placeholders[i].text = str(row[placeholders[i]])
            except KeyError:
                print(f"O slide não tem um placeholder com o índice {i}")

    # Salva a apresentação
    prs.save('nova_promocao_relampago.pptx')

    # Carrega os dados salvos das últimas etiquetas geradas e remove os valores nulos
    path = \
        (
            'https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUeYcOUktnMQ2hknQF7yO7VFGuCO6APMeWCKXk03'
            'Y1ZycTIJ''oAWATed9_PlDoAnnY/pub?gid=1990211285&single=true&output=csv')
    df_salvos = pd.read_csv(path)
    # Remove os valores nulos
    df_salvos = df_salvos.dropna()

    # Altera o tipo de coluna
    df_salvos['Cod Interno'] = df_salvos['Cod Interno'].astype(int)
    df_salvos['Preço DE'] = df_salvos['Preço DE'].str.replace('.', '')
    df_salvos['Preço DE'] = df_salvos['Preço DE'].str.replace(',', '.').astype(float)
    df_salvos['Preço PIX'] = df_salvos['Preço PIX'].str.replace('.', '')
    df_salvos['Preço PIX'] = df_salvos['Preço PIX'].str.replace(',', '.').astype(float)
    df_salvos['Preço Parcelado'] = df_salvos['Preço Parcelado'].str.replace('.', '')
    df_salvos['Preço Parcelado'] = df_salvos['Preço Parcelado'].str.replace(',', '.').astype(float)

    # Primeiro, vamos definir 'Cod Interno' como o índice para ambos os dataframes
    df_salvos.set_index('Cod Interno', inplace=True)
    produtos.set_index('Cod Interno', inplace=True)

    # Agora, vamos remover as linhas em df_salvos que estão em produtos
    df_salvos = df_salvos.loc[~df_salvos.index.isin(produtos.index)]

    # Finalmente, vamos concatenar df_salvos e produtos
    df_final = pd.concat([df_salvos, produtos])

    # Resetando o índice
    df_final.reset_index(inplace=True)

    # Defina o escopo de acesso. Isso especifica quais APIs do Google o código terá acesso.
    scope = ['https://spreadsheets.google.com/feeds', 'https://www.googleapis.com/auth/drive']

    # Use as credenciais do arquivo json para autenticar
    # 'credentials_serv.json' é o arquivo de credenciais da conta de serviço que você baixou do Google Cloud Console.
    creds = ServiceAccountCredentials.from_json_keyfile_name(json_path, scope)

    # Autorize usando as credenciais e inicialize o cliente gspread
    client = gspread.authorize(creds)

    # Abra a planilha do Google usando sua chave única e obtenha a primeira folha
    sheet = client.open_by_key('10emhAinbHEyKEEeQADtFx4yGYe9-JOFrTAzB0_GeLfY').worksheet('Salvos.Promoção Relâmpago')

    # Limpe toda a planilha antes de escrever os dados
    sheet.clear()

    # Atualize a planilha com os dados do DataFrame
    # df_final é o seu DataFrame pandas.
    # Primeiro, convertemos os nomes das colunas do DataFrame em uma lista.
    # Em seguida, convertemos cada linha do DataFrame em uma lista.
    # Finalmente, usamos o método update para escrever essas listas na planilha do Google Sheets.
    sheet.update([df_final.columns.values.tolist()] + df_final.values.tolist())


def generica_amarela():
    # Verifica se estamos executando o código como um script congelado (empacotado pelo PyInstaller)
    if getattr(sys, 'frozen', False):
        # Se estivermos executando como um script congelado, o arquivo .pptx estará no diretório do executável
        pptx_path = os.path.join(sys._MEIPASS, 'generica_amarela.pptx')
        json_path = os.path.join(sys._MEIPASS, 'credentials_serv.json')
    else:
        # Se não estivermos executando como um script congelado, o arquivo .pptx estará no diretório de trabalho atual
        pptx_path = 'generica_amarela.pptx'
        json_path = 'credentials_serv.json'

    # Agora você pode abrir o arquivo .pptx usando o caminho em pptx_path
    with open(pptx_path, 'rb') as f:
        pass  # Faça algo com o arquivo

    # Acessa os dados de uma Planilha Google publicadas em CSV
    path = ('https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUeYcOUktnMQ2hknQF7yO7VFGuCO6APMeWCKXk03'
            'Y1ZycTIJoAWATed9_PlDoAnnY/pub?gid=1067594308&single=true&output=csv')
    produtos = pd.read_csv(path)

    # Remove os valores nulos
    produtos = produtos.dropna()

    # Alterna os tipos de dados de cada coluna
    produtos['Cod Interno'] = produtos['Cod Interno'].astype(int)
    produtos['Preço DE'] = produtos['Preço DE'].str.replace('.', '')
    produtos['Preço DE'] = produtos['Preço DE'].str.replace(',', '.').astype(float)
    produtos['Preço PIX'] = produtos['Preço PIX'].str.replace('.', '')
    produtos['Preço PIX'] = produtos['Preço PIX'].str.replace(',', '.').astype(float)
    produtos['Preço Parcelado'] = produtos['Preço Parcelado'].str.replace('.', '')
    produtos['Preço Parcelado'] = produtos['Preço Parcelado'].str.replace(',', '.').astype(float)
    produtos['PREÇO GARANTIA'] = produtos['PREÇO GARANTIA'].str.replace('.', '')
    produtos['PREÇO GARANTIA'] = produtos['PREÇO GARANTIA'].str.replace(',', '.').astype(float)

    # Cria um novo dataframe para armazenar os dados formatados
    produtos_formatados = (
        pd.DataFrame(columns=['Nome do Produto', 'Cod Interno', 'Preço DE', 'Preço PIX', 'Preço Parcelado', 'ITEM 1',
                              'ITEM 2', 'ITEM 3', 'ITEM 4', 'PREÇO GARANTIA']))

    # Convertendo e formatando as colunas conforme especificado
    produtos_formatados['Cod Interno'] = produtos['Cod Interno']
    produtos_formatados['Nome do Produto'] = produtos['Nome do Produto'].str.upper()
    produtos_formatados['ITEM 1'] = produtos['ITEM 1'].str.upper()
    produtos_formatados['ITEM 2'] = produtos['ITEM 2'].str.upper()
    produtos_formatados['ITEM 3'] = produtos['ITEM 3'].str.upper()
    produtos_formatados['ITEM 4'] = produtos['ITEM 4'].str.upper()
    produtos_formatados['Preço DE'] = ("R$ " + produtos['Preço DE'].apply(
        lambda x: "{:,.2f}".format(x).replace('.', '#').replace(',', '.').replace('#', ',')))
    produtos_formatados['Preço PIX'] = ("R$ " + produtos['Preço PIX'].apply(
        lambda x: "{:,.2f}".format(x).replace('.', '#').replace(',', '.').replace('#', ',')))
    produtos_formatados['PREÇO GARANTIA'] = ("+ 12x R$ " + (produtos['PREÇO GARANTIA'] / 12).apply(
        lambda x: "{:,.2f}".format(x).replace('.', '#').replace(',', '.').replace('#', ',')))

    def calcular_preco_parcelado(preco):
        if preco > 800:
            parcelas = 10
        else:
            parcelas = math.floor(preco / 80)
        return ("Ou " +
                str(parcelas) + "x DE R$ " + "{:,.2f}"
                .format(preco / parcelas).replace(
                    '.', '#').replace(',', '.').replace('#', ','))

    produtos_formatados['Preço Parcelado'] = produtos['Preço Parcelado'].apply(calcular_preco_parcelado)
    print('Sucesso!')

    # Carrega a apresentação
    prs = Presentation(pptx_path)  # Aqui está a mudança

    # Adiciona um novo slide com o layout '2
    slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '

    # Para cada linha do dataframe
    for index, row in produtos_formatados.iterrows():
        # Adiciona um novo slide com o layout '2
        slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '
        slide = prs.slides.add_slide(slide_layout)

        # Mapeia cada coluna para um placeholder específico
        placeholders = {10: 'Nome do Produto', 11: 'Cod Interno', 13: 'Preço DE', 14: 'Preço PIX',
                        15: 'Preço Parcelado',
                        16: 'ITEM 1', 17: 'ITEM 2', 19: 'ITEM 3', 20: 'ITEM 4', 22: 'PREÇO GARANTIA'}

        # Para cada placeholder no slide
        for i in placeholders:
            try:
                # Adiciona o texto do dataframe ao placeholder
                slide.placeholders[i].text = str(row[placeholders[i]])
            except KeyError:
                print(f"O slide não tem um placeholder com o índice {i}")

    # Salva a apresentação
    prs.save('nova_generica_amarela.pptx')

    # Carrega os dados salvos das últimas etiquetas geradas e remove os valores nulos
    path = ('https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUeYcOUktnMQ2hknQF7yO7VFGuCO6APMeWCKXk03Y1Z'
            'ycTIJoAWATed9_PlDoAnnY/pub?gid=1496794293&single=true&output=csv')
    df_salvos = pd.read_csv(path)

    # vRemove os valores nulos
    df_salvos = df_salvos.dropna()
    # Mantém apenas as colunas desejadas
    colunas_desejadas = ['Cod Interno', 'Nome do Produto', 'ITEM 1', 'ITEM 2', 'ITEM 3', 'ITEM 4']

    df_salvos = df_salvos[colunas_desejadas]
    df_salvos['Cod Interno'] = df_salvos['Cod Interno'].astype(int)
    produtos_novos = produtos[colunas_desejadas]
    produtos_novos.loc[:, 'Cod Interno'] = produtos_novos['Cod Interno'].astype(int)

    # Primeiro, vamos definir 'Cod Interno' como o índice para ambos os dataframes
    df_salvos.set_index('Cod Interno', inplace=True)
    produtos_novos.set_index('Cod Interno', inplace=True)

    # Agora, vamos remover as linhas em df_salvos que estão em produtos
    df_salvos = df_salvos.loc[~df_salvos.index.isin(produtos_novos.index)]

    # Finalmente, vamos concatenar df_salvos e produtos
    df_final = pd.concat([df_salvos, produtos_novos])

    # Resetando o índice
    df_final.reset_index(inplace=True)

    # Defina o escopo de acesso. Isso especifica quais APIs do Google o código terá acesso.
    scope = ['https://spreadsheets.google.com/feeds', 'https://www.googleapis.com/auth/drive']

    # Use as credenciais do arquivo json para autenticar
    # 'credentials_serv.json' é o arquivo de credenciais da conta de serviço que você baixou do Google Cloud Console.
    creds = ServiceAccountCredentials.from_json_keyfile_name(json_path, scope)

    # Autorize usando as credenciais e inicialize o cliente gspread
    client = gspread.authorize(creds)

    # Abra a planilha do Google usando sua chave única e obtenha a primeira folha
    sheet = client.open_by_key('10emhAinbHEyKEEeQADtFx4yGYe9-JOFrTAzB0_GeLfY').worksheet('Salvos.Genérica Branca')

    # Limpe toda a planilha antes de escrever os dados
    sheet.clear()

    # Atualize a planilha com os dados do DataFrame
    # df_final é o seu DataFrame pandas.
    # Primeiro, convertemos os nomes das colunas do DataFrame em uma lista.
    # Em seguida, convertemos cada linha do DataFrame em uma lista.
    # Finalmente, usamos o método update para escrever essas listas na planilha do Google Sheets.
    sheet.update([df_final.columns.values.tolist()] + df_final.values.tolist())


def generica_branca():

    # Verifica se estamos executando o código como um script congelado (empacotado pelo PyInstaller)
    if getattr(sys, 'frozen', False):
        # Se estivermos executando como um script congelado, o arquivo .pptx estará no diretório do executável
        pptx_path = os.path.join(sys._MEIPASS, 'generica_branca.pptx')
        json_path = os.path.join(sys._MEIPASS, 'credentials_serv.json')
    else:
        # Se não estivermos executando como um script congelado, o arquivo .pptx estará no diretório de trabalho atual
        pptx_path = 'generica_branca.pptx'
        json_path = 'credentials_serv.json'

    # Agora você pode abrir o arquivo .pptx usando o caminho em pptx_path
    with open(pptx_path, 'rb') as f:
        pass  # Faça algo com o arquivo

    # Acessa os dados de uma Planilha Google publicadas em CSV
    path = ('https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUe'
            'YcOUktnMQ2hknQF7yO7VFGuCO6APMeWCKXk03Y1ZycTIJoAWATed9_PlDoAnnY/pub?gid=449249962&single=true&output=csv')
    produtos = pd.read_csv(path)
    # Remove os valores nulos
    produtos = produtos.dropna()

    # Altera o tipo de dados
    produtos['Cod Interno'] = produtos['Cod Interno'].astype(int)
    produtos['Preço Por'] = produtos['Preço Por'].str.replace('.', '')
    produtos['Preço Por'] = produtos['Preço Por'].str.replace(',', '.').astype(float)
    produtos['PREÇO GARANTIA'] = produtos['PREÇO GARANTIA'].str.replace('.', '')
    produtos['PREÇO GARANTIA'] = produtos['PREÇO GARANTIA'].str.replace(',', '.').astype(float)

    # Cria um novo dataframe para armazenar os dados formatados
    produtos_formatados = (pd.DataFrame(
        columns=['Cod Interno', 'Nome do Produto', 'Preço Por', 'Preço Parcelado', 'ITEM 1', 'ITEM 2', 'ITEM 3',
                 'ITEM 4', 'PREÇO GARANTIA']))

    # Convertendo e formatando as colunas conforme especificado
    produtos_formatados['Cod Interno'] = produtos['Cod Interno']
    produtos_formatados['Nome do Produto'] = produtos['Nome do Produto'].str.upper()
    produtos_formatados['ITEM 1'] = produtos['ITEM 1'].str.upper()
    produtos_formatados['ITEM 2'] = produtos['ITEM 2'].str.upper()
    produtos_formatados['ITEM 3'] = produtos['ITEM 3'].str.upper()
    produtos_formatados['ITEM 4'] = produtos['ITEM 4'].str.upper()
    produtos_formatados['Preço Por'] = (
            "R$ " + produtos['Preço Por'].apply(lambda x: "{:,.2f}".format(x).replace('.', '#').
                                                replace(',', '.').replace('#', ',')))

    produtos_formatados['PREÇO GARANTIA'] = (
            "+ 12x R$ " +
            (produtos['PREÇO GARANTIA'] /
             12).apply(lambda x: "{:,.2f}".format(x).replace(
                '.', '#').replace(',', '.').replace('#', ',')))

    def calcular_preco_parcelado(preco):
        if preco > 800:
            parcelas = 10
        else:
            parcelas = math.floor(preco / 80)
        return ("Ou " + str(parcelas) +
                "x DE R$ " + "{:,.2f}".format(preco / parcelas).replace(
                    '.', '#').replace(',', '.').replace('#', ','))

    produtos_formatados['Preço Parcelado'] = produtos['Preço Por'].apply(calcular_preco_parcelado)

    print('Sucesso!')

    # Carrega a apresentação
    prs = Presentation(pptx_path)  # Aqui está a mudança

    # Adiciona um novo slide com o layout '
    slide_layout = prs.slide_layouts[0]  # Substitua 1 pelo índice do layout '

    # Para cada linha do dataframe
    for index, row in produtos_formatados.iterrows():
        # Adiciona um novo slide com o layout '
        slide_layout = prs.slide_layouts[0]  #
        slide = prs.slides.add_slide(slide_layout)

        # Mapeia cada coluna para um placeholder específico
        placeholders = {10: 'Nome do Produto', 11: 'Cod Interno', 14: 'Preço Por', 15: 'Preço Parcelado', 16: 'ITEM 1',
                        17: 'ITEM 2', 19: 'ITEM 3', 20: 'ITEM 4', 22: 'PREÇO GARANTIA'}

        # Para cada placeholder no slide
        for i in placeholders:
            try:
                # Adiciona o texto do dataframe ao placeholder
                slide.placeholders[i].text = str(row[placeholders[i]])
            except KeyError:
                print(f"O slide não tem um placeholder com o índice {i}")

    # Salva a apresentação
    prs.save('nova_generica_branca.pptx')

    # Carrega os dados salvos das últimas etiquetas geradas e remove os valores nulos
    path = ('https://docs.google.com/spreadsheets/d/e/2PACX-1vTIjvmk4POoaqDWUeYcOUktnMQ2hknQF7yO7VFGuCO6APMeWCKXk03Y1'
            'ZycTIJoAWATed9_PlDoAnnY/pub?gid=1496794293&single=true&output=csv')
    df_salvos = pd.read_csv(path)

    # vRemove os valores nulos
    df_salvos = df_salvos.dropna()
    # Mantém apenas as colunas desejadas
    colunas_desejadas = ['Cod Interno', 'Nome do Produto', 'ITEM 1', 'ITEM 2', 'ITEM 3', 'ITEM 4']

    df_salvos = df_salvos[colunas_desejadas]
    df_salvos['Cod Interno'] = df_salvos['Cod Interno'].astype(int)
    produtos_novos = produtos[colunas_desejadas]
    produtos_novos.loc[:, 'Cod Interno'] = produtos_novos['Cod Interno'].astype(int)

    # Primeiro, vamos definir 'Cod Interno' como o índice para ambos os dataframes
    df_salvos.set_index('Cod Interno', inplace=True)
    produtos_novos.set_index('Cod Interno', inplace=True)

    # Agora, vamos remover as linhas em df_salvos que estão em produtos
    df_salvos = df_salvos.loc[~df_salvos.index.isin(produtos_novos.index)]

    # Finalmente, vamos concatenar df_salvos e produtos
    df_final = pd.concat([df_salvos, produtos_novos])

    # Resetando o índice
    df_final.reset_index(inplace=True)

    # Defina o escopo de acesso. Isso especifica quais APIs do Google o código terá acesso.
    scope = ['https://spreadsheets.google.com/feeds', 'https://www.googleapis.com/auth/drive']

    # Use as credenciais do arquivo json para autenticar
    # 'credentials_serv.json' é o arquivo de credenciais da conta de serviço que você baixou do Google Cloud Console.
    creds = ServiceAccountCredentials.from_json_keyfile_name(json_path, scope)

    # Autorize usando as credenciais e inicialize o cliente gspread
    client = gspread.authorize(creds)

    # Abra a planilha do Google usando sua chave única e obtenha a primeira folha
    sheet = client.open_by_key('10emhAinbHEyKEEeQADtFx4yGYe9-JOFrTAzB0_GeLfY').worksheet('Salvos.Genérica Branca')

    # Limpe toda a planilha antes de escrever os dados
    sheet.clear()

    # Atualize a planilha com os dados do DataFrame
    # df_final é o seu DataFrame pandas.
    # Primeiro, convertemos os nomes das colunas do DataFrame em uma lista.
    # Em seguida, convertemos cada linha do DataFrame em uma lista.
    # Finalmente, usamos o método update para escrever essas listas na planilha do Google Sheets.
    sheet.update([df_final.columns.values.tolist()] + df_final.values.tolist())

def abrir_url():
    webbrowser.open(
        'https://docs.google.com/spreadsheets/d/10emhAinbHEyKEEeQADtFx4yGYe9-JOFrTAzB0_GeLfY/edit#gid=140721476')

# Esta função será chamada quando o botão for pressionado
def on_button_press():
    funcs = [var1.get(), var2.get(), var3.get(), var4.get(), var5.get(), var6.get(), var7.get()]
    if not any(funcs):
        messagebox.showinfo("Erro", "Escolha uma opção")
    else:
        label_aguarde.pack()  # Adiciona a mensagem "Aguarde..." na janela
        root.update()  # Atualiza a janela para exibir a mensagem
        if var1.get():
            celular_amarela()
        if var2.get():
            celular_branca()
        if var3.get():
            notebook_amarela()
        if var4.get():
            notebook_branca()
        if var5.get():
            relampago()
        if var6.get():
            generica_amarela()
        if var7.get():
            generica_branca()
        label_aguarde.pack_forget()  # Remove a mensagem "Aguarde..." da janela

root = tk.Tk()
# Cria a janela principal
root.geometry('400x500')  # Define o tamanho da janela
root.title('GERADOR')
root.configure(bg='#ddebf7')

# Adiciona um título em negrito com fundo azul e fonte branca
titulo = tk.Label(root, text='Gerador de Etiquetas', bg='#0b5394', fg='white', font=('Arial', 12, 'bold'))
titulo.pack(fill=tk.X)  # Faz o título ocupar toda a largura da janela

# Adiciona um título em negrito com fundo azul e fonte branca
titulo = tk.Label(root, text='Selecione uma ou mais etiquetas que deseja gerar.', bg='#0b5394', fg='white',
                  font=('Arial', 9, 'bold'))
titulo.pack(fill=tk.X)  # Faz o título ocupar toda a largura da janela

# Adiciona um título em negrito com fundo azul e fonte branca
titulo = tk.Label(root, text='Compatível com Windows 10 ou superior.',
                  bg='#0b5394', fg='white', font=('Arial', 9, 'bold'))
titulo.pack(fill=tk.X)  # Faz o título ocupar toda a largura da janela

var1 = tk.BooleanVar()
var2 = tk.BooleanVar()
var3 = tk.BooleanVar()
var4 = tk.BooleanVar()
var5 = tk.BooleanVar()
var6 = tk.BooleanVar()
var7 = tk.BooleanVar()

check1 = tk.Checkbutton(root, text="Celular Amarela", variable=var1, font=('Arial', 14, 'bold'),
                        bg='#ddebf7', fg='#0b5394', anchor=tk.W)
check2 = tk.Checkbutton(root, text="Celular Branca", variable=var2, font=('Arial', 14, 'bold'),
                        bg='#ddebf7', fg='#0b5394', anchor=tk.W)
check3 = tk.Checkbutton(root, text="Notebook Amarela", variable=var3, font=('Arial', 14, 'bold'),
                        bg='#ddebf7', fg='#0b5394', anchor=tk.W)
check4 = tk.Checkbutton(root, text="Notebook Branca", variable=var4, font=('Arial', 14, 'bold'),
                        bg='#ddebf7', fg='#0b5394', anchor=tk.W)
check5 = tk.Checkbutton(root, text="Relâmpago", variable=var5, font=('Arial', 14, 'bold'),
                        bg='#ddebf7', fg='#0b5394', anchor=tk.W)
check6 = tk.Checkbutton(root, text="Genérica Amarela", variable=var6, font=('Arial', 14, 'bold'),
                        bg='#ddebf7', fg='#0b5394', anchor=tk.W)
check7 = tk.Checkbutton(root, text="Genérica Branca", variable=var7, font=('Arial', 14, 'bold'), bg='#ddebf7',
                        fg='#0b5394', anchor=tk.W)
check1.pack(anchor=tk.W, padx=100)
check2.pack(anchor=tk.W, padx=100)
check3.pack(anchor=tk.W, padx=100)
check4.pack(anchor=tk.W, padx=100)
check5.pack(anchor=tk.W, padx=100)
check6.pack(anchor=tk.W, padx=100)
check7.pack(anchor=tk.W, padx=100)

button = tk.Button(root, text="Gerar", command=on_button_press)
button.config(bg='#ddebf7', fg='#0b5394', font=('Arial', 12, 'bold'))

check1.pack()
check2.pack()
check3.pack()
check4.pack()
check5.pack()
check6.pack()
check7.pack()

button.pack()

# Cria o botão que abre a URL
botao_url = tk.Button(root, text='Lista de produtos', command=abrir_url)
botao_url.config(bg='#ddebf7', fg='#0b5394', font=('Arial', 12, 'bold'))
botao_url.pack()

# Label para a mensagem de aguarde
label_aguarde = tk.Label(root, text='Aguarde...', bg='#ddebf7', fg='#0b5394', font=('Arial', 12, 'bold'))

# Cria um rodapé
frame_rodape = tk.Frame(root, height=50, bg='#0b5394')
frame_rodape.pack(fill=tk.X, side=tk.BOTTOM)
label_rodape_info = tk.Label(frame_rodape, text='Versão 1.0 Maio 2024 Autor: Robertson Tompson',
                             bg='#0b5394', fg='white', font=('Arial', 10, 'bold'))
label_rodape_info.pack(side=tk.LEFT, padx=10)

root.mainloop()